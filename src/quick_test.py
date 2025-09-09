#!/usr/bin/env python3
"""
快速测试脚本 - 验证PPT替换功能
"""

import os
import sys
from ppt_processor import PPTProcessor

def test_replacement():
    """测试替换功能"""
    print("=" * 50)
    print("PPT替换功能测试")
    print("=" * 50)
    
    # 检查模板文件
    template_path = "file/ces.pptx"
    print(f"检查模板文件: {template_path}")
    
    if not os.path.exists(template_path):
        print(f"❌ 模板文件不存在: {template_path}")
        print("请确保提供的PPT模板文件存在")
        return False
    
    print(f"✓ 模板文件存在")
    
    try:
        # 创建处理器
        processor = PPTProcessor(template_path)
        print("✓ 创建PPT处理器成功")
        
        # 加载模板
        if not processor.load_template():
            print("❌ 加载模板失败")
            return False
        
        print("✓ 模板加载成功")
        
        # 查找占位符
        placeholders = processor.find_placeholders()
        print(f"✓ 找到 {len(placeholders)} 个占位符:")
        
        for slide_idx, placeholder in placeholders:
            print(f"  - 幻灯片 {slide_idx + 1}: {placeholder}")
        
        if len(placeholders) == 0:
            print("❌ 未找到任何占位符，请检查PPT模板是否包含{{analysis_text1}}或{{analysis_text2}}")
            return False
        
        # 测试替换
        output_path = "test_output.pptx"
        print(f"\n开始处理PPT，输出文件: {output_path}")
        
        if processor.process(output_path):
            print("✓ PPT处理成功")
            
            if os.path.exists(output_path):
                file_size = os.path.getsize(output_path)
                print(f"✓ 输出文件已生成: {output_path}")
                print(f"✓ 文件大小: {file_size / 1024:.1f} KB")
                
                # 验证输出文件
                try:
                    from pptx import Presentation
                    test_ppt = Presentation(output_path)
                    print(f"✓ 输出文件格式正确，包含 {len(test_ppt.slides)} 张幻灯片")
                    
                    # 检查是否还有占位符
                    remaining_placeholders = []
                    for slide_idx, slide in enumerate(test_ppt.slides):
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame"):
                                text = shape.text_frame.text
                                import re
                                matches = re.findall(r'\{\{[^}]+\}\}', text)
                                for match in matches:
                                    remaining_placeholders.append((slide_idx, match))
                    
                    if remaining_placeholders:
                        print(f"⚠️  仍有 {len(remaining_placeholders)} 个占位符未被替换:")
                        for slide_idx, placeholder in remaining_placeholders:
                            print(f"  - 幻灯片 {slide_idx + 1}: {placeholder}")
                    else:
                        print("✓ 所有占位符已成功替换")
                    
                except Exception as e:
                    print(f"❌ 验证输出文件时出错: {e}")
                    return False
                    
            else:
                print("❌ 输出文件未生成")
                return False
        else:
            print("❌ PPT处理失败")
            return False
    
    except Exception as e:
        print(f"❌ 测试过程中出现异常: {e}")
        import traceback
        traceback.print_exc()
        return False
    
    print("\n" + "=" * 50)
    print("测试完成")
    print("=" * 50)
    return True

if __name__ == "__main__":
    # 检查依赖
    try:
        import pptx
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
    except ImportError as e:
        print(f"❌ 缺少依赖库: {e}")
        print("请运行: pip install python-pptx")
        sys.exit(1)
    
    # 运行测试
    if test_replacement():
        print("🎉 测试成功！")
        sys.exit(0)
    else:
        print("❌ 测试失败！")
        sys.exit(1)