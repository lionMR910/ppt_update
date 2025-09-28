#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
增强的PPT处理器 - 集成KPI参数替换功能

扩展现有的PPTProcessor类，增加KPI参数处理功能
"""

import re
import os
import logging
from typing import Dict, List, Tuple, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# 导入现有模块
try:
    from .ppt_processor import PPTProcessor
    from .kpi_replacer import KpiReplacer
    from .database_config import DatabaseManager, DatabaseConfig
except ImportError:
    from ppt_processor import PPTProcessor
    from kpi_replacer import KpiReplacer
    from database_config import DatabaseManager, DatabaseConfig


class EnhancedPPTProcessor(PPTProcessor):
    """增强的PPT处理器，支持KPI参数替换"""
    
    def __init__(self, template_path: str, db_manager: DatabaseManager = None, analysis_id: int = 1, sql_ids: list = None):
        """
        初始化增强的PPT处理器
        
        Args:
            template_path: PPT模板文件路径
            db_manager: 数据库管理器实例，可选
            analysis_id: 分析配置ID，1=全球通"量质构效"分析，2=中高端"量质构效"分析
            sql_ids: 指定的分析任务SQL ID列表，为空则处理所有任务
        """
        super().__init__(template_path)
        self.logger = logging.getLogger(__name__)
        
        # 初始化KPI替换器
        self.kpi_replacer = KpiReplacer(db_manager, analysis_id, sql_ids)
        
        # 当前月份参数
        self.current_month = None
        
        # KPI相关的替换数据
        self.kpi_replacement_data = {}
    
    def set_current_month(self, month: str) -> bool:
        """
        设置当前月份参数
        
        Args:
            month: 月份参数，格式为YYYYMM
            
        Returns:
            设置是否成功
        """
        try:
            # 验证月份格式
            if not re.match(r'^\d{6}$', month):
                raise ValueError("月份格式错误，应为YYYYMM格式")
            
            year = int(month[:4])
            month_num = int(month[4:])
            
            if month_num < 1 or month_num > 12:
                raise ValueError("月份数值错误，应在01-12之间")
            
            if year < 2020 or year > 2030:
                raise ValueError("年份超出合理范围（2020-2030）")
            
            self.current_month = month
            self.logger.info(f"设置当前月份: {month}")
            return True
            
        except Exception as e:
            self.logger.error(f"设置月份失败: {e}")
            return False
    
    def process_kpi_placeholders(self) -> bool:
        """
        处理KPI占位符，获取数据库数据并准备替换
        
        Returns:
            处理是否成功
        """
        if not self.current_month:
            self.logger.error("未设置当前月份参数")
            return False
        
        try:
            # 查找PPT中的所有KPI占位符
            kpi_placeholders = self._find_kpi_placeholders()
            
            if not kpi_placeholders:
                self.logger.info("PPT中未找到KPI占位符")
                return True
            
            self.logger.info(f"找到 {len(kpi_placeholders)} 个KPI占位符")
            
            # 获取KPI数据
            kpi_values = self.kpi_replacer.get_kpi_values(self.current_month)
            
            # 准备替换数据
            for placeholder, sql_id, column_index in kpi_placeholders:
                key = (sql_id, column_index)
                
                if key in kpi_values:
                    value = kpi_values[key]
                    self.kpi_replacement_data[placeholder] = str(value)
                    self.logger.info(f"准备替换: {placeholder} -> {value}")
                else:
                    self.logger.warning(f"未找到KPI数据: {placeholder}")
                    # 使用默认值或保持原样
                    self.kpi_replacement_data[placeholder] = placeholder
            
            return True
            
        except Exception as e:
            self.logger.error(f"处理KPI占位符失败: {e}")
            return False
    
    def _find_kpi_placeholders(self) -> List[Tuple[str, int, int]]:
        """
        在PPT中查找所有KPI占位符
        
        Returns:
            KPI占位符列表，每个元素为(完整占位符, sql_id, 列索引)
        """
        placeholders = []
        
        if not self.presentation:
            return placeholders
        
        # 遍历所有幻灯片
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    if text:
                        # 提取KPI占位符
                        found_placeholders = self.kpi_replacer.extract_kpi_placeholders(text)
                        placeholders.extend(found_placeholders)
        
        # 去重
        unique_placeholders = list(set(placeholders))
        return unique_placeholders
    
    def replace_kpi_placeholders(self) -> Dict[str, int]:
        """
        替换PPT中的KPI占位符，保持原有格式
        
        Returns:
            替换结果统计
        """
        results = {"success": 0, "failed": 0, "not_found": 0}
        
        if not self.kpi_replacement_data:
            self.logger.warning("没有准备好的KPI替换数据")
            return results
        
        # 遍历所有幻灯片进行替换
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # 使用精确的文本运行块替换，保持原有格式
                    replacements_made = self._replace_text_in_runs(shape.text_frame)
                    if replacements_made > 0:
                        results["success"] += replacements_made
        
        self.logger.info(f"KPI替换完成: 成功{results['success']}个，失败{results['failed']}个")
        return results
    
    def _replace_text_in_runs(self, text_frame) -> int:
        """
        在文本运行块中进行精确替换，处理跨运行块的占位符
        
        Args:
            text_frame: PPT文本框架
            
        Returns:
            替换的占位符数量
        """
        total_replacements = 0
        
        try:
            # 遍历所有段落
            for paragraph in text_frame.paragraphs:
                # 收集段落中所有文本和运行块信息
                full_text = ""
                run_info = []  # [(start_pos, end_pos, run_obj), ...]
                
                for run in paragraph.runs:
                    start_pos = len(full_text)
                    run_text = run.text if run.text else ""
                    full_text += run_text
                    end_pos = len(full_text)
                    run_info.append((start_pos, end_pos, run))
                
                if not full_text:
                    continue
                
                # 在完整文本中查找和替换占位符
                modified_text = full_text
                paragraph_replacements = 0
                
                for placeholder, value in self.kpi_replacement_data.items():
                    if placeholder in modified_text:
                        old_count = modified_text.count(placeholder)
                        modified_text = modified_text.replace(placeholder, str(value))
                        new_count = modified_text.count(placeholder)
                        replacements_made = old_count - new_count
                        paragraph_replacements += replacements_made
                        self.logger.info(f"在段落中替换: {placeholder} -> {value} (替换了{replacements_made}次)")
                
                # 如果有替换，重新构建段落的运行块
                if paragraph_replacements > 0:
                    self._reconstruct_paragraph_runs(paragraph, modified_text, run_info)
                    total_replacements += paragraph_replacements
                else:
                    # 检查单个运行块中的替换（向后兼容）
                    for run in paragraph.runs:
                        if run.text:
                            original_text = run.text
                            new_text = original_text
                            run_replacements = 0
                            
                            for placeholder, value in self.kpi_replacement_data.items():
                                if placeholder in new_text:
                                    old_count = new_text.count(placeholder)
                                    new_text = new_text.replace(placeholder, str(value))
                                    new_count = new_text.count(placeholder)
                                    run_replacements += (old_count - new_count)
                            
                            if run_replacements > 0:
                                run.text = new_text
                                total_replacements += run_replacements
                                self.logger.debug(f"单个运行块替换: {run_replacements} 个占位符")
                            
        except Exception as e:
            self.logger.error(f"运行块替换失败: {e}")
            
        return total_replacements
    
    def _reconstruct_paragraph_runs(self, paragraph, new_text: str, original_run_info: list):
        """
        根据新文本重新构建段落的运行块，尽量保持格式
        
        Args:
            paragraph: PPT段落对象
            new_text: 新的文本内容
            original_run_info: 原始运行块信息 [(start_pos, end_pos, run_obj), ...]
        """
        try:
            # 保存第一个运行块的格式
            original_format = None
            if original_run_info and len(original_run_info) > 0:
                first_run = original_run_info[0][2]
                original_format = self._extract_run_format(first_run)
            
            # 清空现有运行块
            paragraph.clear()
            
            # 创建新的运行块
            if new_text:
                new_run = paragraph.add_run()
                new_run.text = new_text
                
                # 应用原始格式
                if original_format:
                    self._apply_run_format(new_run, original_format)
                    
        except Exception as e:
            self.logger.error(f"重构段落运行块失败: {e}")
    
    def _extract_run_format(self, run):
        """
        提取运行块的格式信息
        
        Args:
            run: 原始运行块
            
        Returns:
            格式信息字典
        """
        try:
            return {
                'font_name': getattr(run.font, 'name', None),
                'font_size': getattr(run.font, 'size', None),
                'bold': getattr(run.font, 'bold', None),
                'italic': getattr(run.font, 'italic', None),
                'underline': getattr(run.font, 'underline', None),
                'color': getattr(run.font.color, 'rgb', None) if hasattr(run.font, 'color') else None
            }
        except Exception as e:
            self.logger.debug(f"提取格式信息失败: {e}")
            return {}
    
    def _apply_run_format(self, run, format_info: dict):
        """
        应用格式信息到运行块
        
        Args:
            run: 目标运行块
            format_info: 格式信息字典
        """
        try:
            if format_info.get('font_name'):
                run.font.name = format_info['font_name']
            if format_info.get('font_size'):
                run.font.size = format_info['font_size']
            if format_info.get('bold') is not None:
                run.font.bold = format_info['bold']
            if format_info.get('italic') is not None:
                run.font.italic = format_info['italic']
            if format_info.get('underline') is not None:
                run.font.underline = format_info['underline']
            if format_info.get('color'):
                run.font.color.rgb = format_info['color']
        except Exception as e:
            self.logger.debug(f"应用格式失败: {e}")
    
    def _replace_text_with_format(self, shape, new_text: str) -> bool:
        """
        替换形状中的文本，保持原有格式不变
        
        Args:
            shape: PPT形状对象
            new_text: 新文本内容
            
        Returns:
            替换是否成功
        """
        try:
            # 直接替换文本内容，保持原有格式
            shape.text_frame.text = new_text
            return True
            
        except Exception as e:
            self.logger.error(f"文本替换失败: {e}")
            return False
    
    def _apply_text_formatting(self, paragraph, text: str):
        """
        应用文本格式化（从父类复制的方法）
        
        Args:
            paragraph: PPT段落对象
            text: 要格式化的文本
        """
        # 清空段落
        paragraph.clear()
        
        # 重新设置段落行间距
        if hasattr(self, 'format_config') and 'line_spacing' in self.format_config:
            paragraph.line_spacing = self.format_config['line_spacing']
        
        # 定义模式
        colon_pattern = r"([\u4e00-\u9fff]{2,15}?)："  # 简化的冒号前文字模式
        number_pattern = r"(\d+(?:\.\d+)?%?)"  # 数字和百分号
        
        # 收集所有需要特殊格式的位置
        special_ranges = []
        
        # 收集冒号前文字
        for match in re.finditer(colon_pattern, text):
            # 只有在句子开头或标点符号后的冒号前文字才加粗
            start_pos = match.start()
            if start_pos == 0 or text[start_pos-1] in '。；\n\r':
                special_ranges.append((match.start(), match.end(), 'bold_red'))
        
        # 收集数字
        for match in re.finditer(number_pattern, text):
            # 检查是否与已有范围重叠
            overlap = False
            for start, end, _ in special_ranges:
                if not (match.end() <= start or match.start() >= end):
                    overlap = True
                    break
            if not overlap:
                special_ranges.append((match.start(), match.end(), 'red_only'))
        
        # 按位置排序
        special_ranges.sort()
        
        # 重新构建段落
        last_pos = 0
        
        for start, end, format_type in special_ranges:
            # 添加普通文本
            if start > last_pos:
                normal_run = paragraph.add_run()
                normal_run.text = text[last_pos:start]
                if hasattr(self, 'format_config'):
                    normal_run.font.name = self.format_config['font_name']
                    normal_run.font.size = Pt(self.format_config['font_size'])
            
            # 添加特殊格式文本
            special_run = paragraph.add_run()
            special_run.text = text[start:end]
            if hasattr(self, 'format_config'):
                special_run.font.name = self.format_config['font_name']
                special_run.font.size = Pt(self.format_config['font_size'])
                
                if format_type == 'bold_red':
                    special_run.font.bold = True
                    special_run.font.color.rgb = RGBColor.from_string(self.format_config['highlight_color'])
                elif format_type == 'red_only':
                    special_run.font.color.rgb = RGBColor.from_string(self.format_config['highlight_color'])
            
            last_pos = end
        
        # 添加剩余的普通文本
        if last_pos < len(text):
            final_run = paragraph.add_run()
            final_run.text = text[last_pos:]
            if hasattr(self, 'format_config'):
                final_run.font.name = self.format_config['font_name']
                final_run.font.size = Pt(self.format_config['font_size'])
    
    def process_complete_replacement(self, month: str) -> Dict[str, int]:
        """
        完整的替换流程：处理传统占位符和KPI占位符
        
        Args:
            month: 当前月份参数，格式为YYYYMM
            
        Returns:
            替换结果统计
        """
        total_results = {"success": 0, "failed": 0, "not_found": 0}
        
        # 1. 设置当前月份
        if not self.set_current_month(month):
            return total_results
        
        # 2. 处理KPI占位符
        if self.process_kpi_placeholders():
            kpi_results = self.replace_kpi_placeholders()
            
            # 合并结果
            for key in total_results:
                total_results[key] += kpi_results.get(key, 0)
        
        # 3. 处理传统占位符（如果有的话）
        if self.replacement_data:
            traditional_results = self.replace_all_placeholders()
            
            # 合并结果
            for key in total_results:
                total_results[key] += traditional_results.get(key, 0)
        
        return total_results
    
    def generate_replacement_report(self, results: Dict[str, int]) -> str:
        """
        生成替换结果报告
        
        Args:
            results: 替换结果统计
            
        Returns:
            报告文本
        """
        report = f"""
PPT参数替换完成报告
===================

替换时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
当前月份: {self.current_month or '未设置'}

替换结果:
- 成功替换: {results.get('success', 0)} 个
- 替换失败: {results.get('failed', 0)} 个
- 未找到占位符: {results.get('not_found', 0)} 个

KPI替换数据:
"""
        
        for placeholder, value in self.kpi_replacement_data.items():
            report += f"- {placeholder}: {value}\n"
        
        if self.replacement_data:
            report += "\n传统替换数据:\n"
            for placeholder, value in self.replacement_data.items():
                # 截断长文本
                display_value = value[:50] + "..." if len(str(value)) > 50 else value
                report += f"- {placeholder}: {display_value}\n"
        
        return report


def main():
    """测试函数"""
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    
    # 测试增强的PPT处理器
    try:
        # 初始化处理器
        processor = EnhancedPPTProcessor("file/ces.pptx")
        
        # 加载模板
        if not processor.load_template():
            print("模板加载失败")
            return
        
        # 执行完整替换流程
        results = processor.process_complete_replacement("202507")
        
        # 生成报告
        report = processor.generate_replacement_report(results)
        print(report)
        
        # 保存结果
        output_path = f"file/ces_kpi_updated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        if processor.save_presentation(output_path):
            print(f"PPT已保存到: {output_path}")
        else:
            print("PPT保存失败")
            
    except Exception as e:
        logging.error(f"测试失败: {e}")


if __name__ == "__main__":
    main()