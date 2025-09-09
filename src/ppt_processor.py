"""
PPT处理模块 - 实现PPT模板文字替换和格式控制
"""

import re
import os
from typing import Dict, List, Tuple
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.text.text import TextFrame, _Run
from config import FORMAT_CONFIG, REPLACEMENT_DATA
from analysis_data_text_order import analysis_data_text
from ppt_data_generator import PPTDataGenerator


class PPTProcessor:
    def __init__(self, template_path: str, sql_ids: list = None):
        """
        初始化PPT处理器
        
        Args:
            template_path: PPT模板文件路径
            sql_ids: 指定的分析任务SQL ID列表，为空则处理所有任务
        """
        self.template_path = template_path
        self.sql_ids = sql_ids
        self.presentation = None
        self.replacement_data = {}  # 动态生成的替换数据
        self.format_config = FORMAT_CONFIG
        
    def load_template(self) -> bool:
        """
        加载PPT模板文件
        
        Returns:
            bool: 加载成功返回True，失败返回False
        """
        try:
            if not os.path.exists(self.template_path):
                raise FileNotFoundError(f"模板文件不存在: {self.template_path}")
            
            self.presentation = Presentation(self.template_path)
            return True
        except Exception as e:
            print(f"加载模板失败: {e}")
            return False
    
    def generate_analysis_data(self, api_key: str, user_input: str, chart_obj: str, conversation_uid: str = None) -> bool:
        """
        生成分析报告数据用于PPT替换
        
        Args:
            api_key: API密钥
            user_input: 用户输入问题
            chart_obj: 图表数据对象
            conversation_uid: 对话ID，可选
            
        Returns:
            bool: 生成成功返回True，失败返回False
        """
        try:
            # 如果没有提供conversation_uid，生成一个
            if conversation_uid is None:
                import random
                conversation_uid = ''.join([str(random.randint(0, 9)) for _ in range(10)])
            
            # 调用分析函数获取分析报告
            analysis_result = analysis_data_text(api_key, user_input, conversation_uid, chart_obj)
            
            if analysis_result:
                # 设置替换数据
                self.replacement_data = {
                    "{{analysis_content}}": analysis_result,
                    "{{analysis_title}}": "数据分析报告",
                    "{{data_source}}": "中国移动数据中心",
                    "{{analysis_date}}": self._get_current_date(),
                }
                return True
            else:
                print("分析报告生成失败")
                return False
                
        except Exception as e:
            print(f"生成分析数据失败: {e}")
            return False
    
    def _get_current_date(self) -> str:
        """获取当前日期字符串"""
        from datetime import datetime
        return datetime.now().strftime("%Y年%m月%d日")
    
    def find_placeholders(self) -> List[Tuple[int, str]]:
        """
        查找PPT中的占位符
        
        Returns:
            List[Tuple[int, str]]: 找到的占位符列表 (幻灯片索引, 占位符文本)
        """
        placeholders = []
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            # 检查文本框
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    # 查找所有{{}}格式的占位符
                    matches = re.findall(r'\{\{[^}]+\}\}', text)
                    for match in matches:
                        placeholders.append((slide_idx, match))
        
        return placeholders
    
    def replace_text_in_shape(self, shape, placeholder: str, replacement_text: str) -> bool:
        """
        在指定形状中替换文本并应用格式
        
        Args:
            shape: PPT形状对象
            placeholder: 占位符文本
            replacement_text: 替换文本
            
        Returns:
            bool: 替换成功返回True
        """
        try:
            if not hasattr(shape, "text_frame"):
                return False
            
            text_frame = shape.text_frame
            
            # 检查整个文本框是否只包含占位符
            full_text = text_frame.text
            if full_text.strip() == placeholder.strip():
                # 清空所有段落
                text_frame.clear()
                
                # 按段落分割替换文本
                paragraphs = replacement_text.strip().split('\n\n')
                
                for i, para_text in enumerate(paragraphs):
                    if para_text.strip():  # 只处理非空段落
                        if i == 0:
                            # 使用第一个段落（已经存在）
                            paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                        else:
                            # 添加新段落
                            paragraph = text_frame.add_paragraph()
                        
                        # 应用格式化文本
                        self._apply_formatted_text(paragraph, para_text.strip())
                
                return True
            else:
                # 如果文本框包含其他内容，进行字符串替换
                for paragraph in text_frame.paragraphs:
                    if placeholder in paragraph.text:
                        # 简单的字符串替换
                        new_text = paragraph.text.replace(placeholder, replacement_text.strip())
                        paragraph.clear()
                        self._apply_formatted_text(paragraph, new_text)
                        return True
            
            return False
        except Exception as e:
            print(f"替换文本时出错: {e}")
            return False
    
    def _apply_formatted_text(self, paragraph, text: str):
        """
        应用格式化文本到段落
        
        Args:
            paragraph: PPT段落对象
            text: 要格式化的文本
        """
        # 设置段落行间距
        if 'line_spacing' in self.format_config:
            paragraph.line_spacing = self.format_config['line_spacing']
        
        # 首先添加默认格式的完整文本
        run = paragraph.add_run()
        run.text = text
        run.font.name = self.format_config['font_name']
        run.font.size = Pt(self.format_config['font_size'])
        
        # 然后处理特殊格式
        self._apply_special_formatting(paragraph, text)
    
    def _apply_special_formatting(self, paragraph, text: str):
        """
        应用特殊格式（红色、加粗）
        这个方法通过字符级别的格式设置来实现特殊格式
        """
        if not paragraph.runs or not paragraph.runs[0].text:
            return
        
        # 获取第一个（也是唯一的）run
        main_run = paragraph.runs[0]
        
        # 清空段落重新构建
        paragraph.clear()
        
        # 重新设置段落行间距
        if 'line_spacing' in self.format_config:
            paragraph.line_spacing = self.format_config['line_spacing']
        
        # 定义模式
        colon_pattern = r"([\u4e00-\u9fff]{2,15}?)："  # 简化的冒号前文字模式
        number_pattern = r"(\d+(?:\.\d+)?%?)"  # 数字和百分号
        
        # 收集所有需要特殊格式的位置
        special_ranges = []
        
        # 收集冒号前文字
        for match in re.finditer(colon_pattern, text):
            # 只有在句子开头或标点符号后的冒号前文字才加粗
            start_pos = match.start()
            if start_pos == 0 or text[start_pos-1] in '。；\n\r':
                special_ranges.append((match.start(), match.end(), 'bold_red'))
        
        # 收集数字
        for match in re.finditer(number_pattern, text):
            # 检查是否与已有范围重叠
            overlap = False
            for start, end, _ in special_ranges:
                if not (match.end() <= start or match.start() >= end):
                    overlap = True
                    break
            if not overlap:
                special_ranges.append((match.start(), match.end(), 'red_only'))
        
        # 按位置排序
        special_ranges.sort()
        
        # 重新构建段落
        last_pos = 0
        
        for start, end, format_type in special_ranges:
            # 添加普通文本
            if start > last_pos:
                normal_run = paragraph.add_run()
                normal_run.text = text[last_pos:start]
                normal_run.font.name = self.format_config['font_name']
                normal_run.font.size = Pt(self.format_config['font_size'])
            
            # 添加特殊格式文本
            special_run = paragraph.add_run()
            special_run.text = text[start:end]
            special_run.font.name = self.format_config['font_name']
            special_run.font.size = Pt(self.format_config['font_size'])
            
            if format_type == 'bold_red':
                special_run.font.bold = True
                special_run.font.color.rgb = RGBColor.from_string(self.format_config['highlight_color'])
            elif format_type == 'red_only':
                special_run.font.color.rgb = RGBColor.from_string(self.format_config['highlight_color'])
            
            last_pos = end
        
        # 添加剩余的普通文本
        if last_pos < len(text):
            final_run = paragraph.add_run()
            final_run.text = text[last_pos:]
            final_run.font.name = self.format_config['font_name']
            final_run.font.size = Pt(self.format_config['font_size'])
    
    def replace_all_placeholders(self) -> Dict[str, int]:
        """
        替换所有占位符
        
        Returns:
            Dict[str, int]: 替换结果统计
        """
        results = {"success": 0, "failed": 0, "not_found": 0}
        
        for placeholder, replacement_text in self.replacement_data.items():
            placeholder_found = False
            replacement_count = 0
            
            # 遍历所有幻灯片
            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        if placeholder in shape.text_frame.text:
                            placeholder_found = True
                            if self.replace_text_in_shape(shape, placeholder, replacement_text):
                                replacement_count += 1
                            else:
                                results["failed"] += 1
            
            if placeholder_found:
                results["success"] += replacement_count
            else:
                results["not_found"] += 1
                print(f"警告: 未找到占位符 {placeholder}")
        
        return results
    
    def save_presentation(self, output_path: str) -> bool:
        """
        保存处理后的PPT文件
        
        Args:
            output_path: 输出文件路径
            
        Returns:
            bool: 保存成功返回True
        """
        try:
            # 确保输出目录存在
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            self.presentation.save(output_path)
            return True
        except Exception as e:
            print(f"保存文件失败: {e}")
            return False
    
    def process(self, output_path: str) -> bool:
        """
        执行完整的处理流程
        
        Args:
            output_path: 输出文件路径
            
        Returns:
            bool: 处理成功返回True
        """
        print("开始处理PPT模板...")
        
        # 加载模板
        if not self.load_template():
            return False
        
        print(f"成功加载模板: {self.template_path}")
        
        # 查找占位符
        placeholders = self.find_placeholders()
        print(f"找到 {len(placeholders)} 个占位符")
        
        # 如果没有替换数据，尝试从数据库动态生成
        if not self.replacement_data:
            print("🔄 正在从数据库生成动态替换数据...")
            # 如果指定了sql_ids，只处理指定任务；否则处理所有任务
            use_all_tasks = self.sql_ids is None
            data_generator = PPTDataGenerator(use_all_tasks=use_all_tasks, sql_ids=self.sql_ids)
            
            if data_generator.connect_database():
                try:
                    dynamic_data = data_generator.generate_replacement_data()
                    if dynamic_data:
                        self.replacement_data = dynamic_data
                        print(f"✅ 成功生成动态替换数据: {list(self.replacement_data.keys())}")
                    else:
                        print("⚠️ 动态数据生成失败，使用静态配置数据")
                        self.replacement_data = REPLACEMENT_DATA
                        print(f"📋 使用配置中的静态替换数据: {list(self.replacement_data.keys())}")
                finally:
                    data_generator.disconnect_database()
            else:
                print("❌ 数据库连接失败，使用静态配置数据")
                self.replacement_data = REPLACEMENT_DATA
                print(f"📋 使用配置中的静态替换数据: {list(self.replacement_data.keys())}")
        
        # 替换占位符
        results = self.replace_all_placeholders()
        print(f"替换结果: 成功 {results['success']}, 失败 {results['failed']}, 未找到 {results['not_found']}")
        
        # 保存文件
        if self.save_presentation(output_path):
            print(f"成功保存到: {output_path}")
            return True
        else:
            return False