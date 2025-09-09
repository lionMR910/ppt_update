#!/usr/bin/env python3
"""
验证替换结果脚本
"""

from pptx import Presentation
import re

def verify_replacement():
    """验证PPT中的文字是否被正确替换"""
    print("验证替换结果...")
    
    try:
        # 加载输出文件
        ppt = Presentation("test_output.pptx")
        
        print(f"PPT包含 {len(ppt.slides)} 张幻灯片")
        
        # 检查每张幻灯片的文字内容
        for slide_idx, slide in enumerate(ppt.slides):
            print(f"\n--- 幻灯片 {slide_idx + 1} ---")
            
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    text = shape.text_frame.text
                    
                    # 检查是否包含我们的替换内容
                    if "沈阳市" in text or "大连市" in text or "营口市" in text:
                        print(f"✓ 形状 {shape_idx + 1} 包含替换后的内容:")
                        print(f"  前50个字符: {text[:50]}...")
                        
                        # 检查是否还有占位符
                        placeholders = re.findall(r'\{\{[^}]+\}\}', text)
                        if placeholders:
                            print(f"  ❌ 仍包含占位符: {placeholders}")
                        else:
                            print(f"  ✓ 无占位符残留")
                    
                    elif "{{" in text:
                        print(f"❌ 形状 {shape_idx + 1} 仍包含未替换的占位符:")
                        print(f"  内容: {text}")
        
        print("\n验证完成")
        return True
        
    except Exception as e:
        print(f"❌ 验证过程中出错: {e}")
        return False

if __name__ == "__main__":
    verify_replacement()