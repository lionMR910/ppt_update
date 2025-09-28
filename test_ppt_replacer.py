#!/usr/bin/env python3
"""
PPT替换程序测试脚本
"""

import os
import sys
import tempfile
import shutil
from ppt_processor import PPTProcessor
from config import REPLACEMENT_DATA


def test_basic_functionality():
    """测试基本功能"""
    print("测试1: 基本功能测试")
    
    # 检查模板文件是否存在
    template_path = "file/ces.pptx"
    if not os.path.exists(template_path):
        print(f"❌ 模板文件不存在: {template_path}")
        print("请确保 file/ces.pptx 文件存在")
        return False
    
    print(f"✓ 模板文件存在: {template_path}")
    
    # 创建处理器
    try:
        processor = PPTProcessor(template_path)
        print("✓ PPT处理器创建成功")
    except Exception as e:
        print(f"❌ PPT处理器创建失败: {e}")
        return False
    
    # 加载模板
    if processor.load_template():
        print("✓ 模板加载成功")
    else:
        print("❌ 模板加载失败")
        return False
    
    # 查找占位符
    placeholders = processor.find_placeholders()
    print(f"✓ 找到 {len(placeholders)} 个占位符")
    
    if placeholders:
        for slide_idx, placeholder in placeholders:
            print(f"  - 幻灯片 {slide_idx + 1}: {placeholder}")
    
    return True


def test_replacement():
    """测试替换功能"""
    print("\n测试2: 替换功能测试")
    
    template_path = "file/ces.pptx"
    if not os.path.exists(template_path):
        print(f"❌ 模板文件不存在: {template_path}")
        return False
    
    # 创建临时输出文件
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
        output_path = temp_file.name
    
    try:
        processor = PPTProcessor(template_path)
        
        # 执行完整处理流程
        if processor.process(output_path):
            print("✓ PPT处理完成")
            
            # 检查输出文件
            if os.path.exists(output_path):
                file_size = os.path.getsize(output_path)
                print(f"✓ 输出文件已生成: {output_path}")
                print(f"✓ 文件大小: {file_size / 1024:.1f} KB")
                
                # 尝试重新加载输出文件以验证格式
                try:
                    from pptx import Presentation
                    test_ppt = Presentation(output_path)
                    print(f"✓ 输出文件格式正确，包含 {len(test_ppt.slides)} 张幻灯片")
                except Exception as e:
                    print(f"❌ 输出文件格式验证失败: {e}")
                    return False
                    
            else:
                print("❌ 输出文件未生成")
                return False
        else:
            print("❌ PPT处理失败")
            return False
    
    except Exception as e:
        print(f"❌ 替换测试失败: {e}")
        return False
    
    finally:
        # 清理临时文件
        if os.path.exists(output_path):
            os.unlink(output_path)
    
    return True


def test_cli_interface():
    """测试命令行界面"""
    print("\n测试3: 命令行界面测试")
    
    template_path = "file/ces.pptx"
    if not os.path.exists(template_path):
        print(f"❌ 模板文件不存在: {template_path}")
        return False
    
    # 测试帮助信息
    print("测试帮助信息...")
    os.system("python ppt_replacer.py --help")
    
    # 测试版本信息
    print("\n测试版本信息...")
    os.system("python ppt_replacer.py --version")
    
    return True


def test_configuration():
    """测试配置文件"""
    print("\n测试4: 配置验证测试")
    
    # 检查替换数据配置
    print("检查替换数据配置...")
    for placeholder, content in REPLACEMENT_DATA.items():
        if not placeholder.startswith('{{') or not placeholder.endswith('}}'):
            print(f"❌ 占位符格式错误: {placeholder}")
            return False
        
        if not content.strip():
            print(f"❌ 替换内容为空: {placeholder}")
            return False
    
    print(f"✓ 配置了 {len(REPLACEMENT_DATA)} 个替换项")
    
    # 检查配置文件导入
    try:
        from config import FORMAT_CONFIG, FILE_CONFIG
        print("✓ 配置文件导入成功")
        print(f"✓ 字体配置: {FORMAT_CONFIG['font_name']} {FORMAT_CONFIG['font_size']}号")
        print(f"✓ 高亮颜色: #{FORMAT_CONFIG['highlight_color']}")
    except Exception as e:
        print(f"❌ 配置文件导入失败: {e}")
        return False
    
    return True


def run_all_tests():
    """运行所有测试"""
    print("=" * 60)
    print("PPT替换程序 - 功能测试")
    print("=" * 60)
    
    tests = [
        test_configuration,
        test_basic_functionality,
        test_replacement,
        test_cli_interface
    ]
    
    passed = 0
    failed = 0
    
    for test_func in tests:
        try:
            if test_func():
                passed += 1
                print("✓ 测试通过\n")
            else:
                failed += 1
                print("❌ 测试失败\n")
        except Exception as e:
            failed += 1
            print(f"❌ 测试异常: {e}\n")
    
    print("=" * 60)
    print(f"测试结果: 通过 {passed}, 失败 {failed}")
    print("=" * 60)
    
    if failed == 0:
        print("🎉 所有测试通过！程序可以正常使用。")
        return True
    else:
        print("⚠️  部分测试失败，请检查程序配置。")
        return False


if __name__ == "__main__":
    # 检查依赖
    try:
        import pptx
        import colorama
    except ImportError as e:
        print(f"❌ 缺少依赖库: {e}")
        print("请运行: pip install -r requirements.txt")
        sys.exit(1)
    
    # 运行测试
    success = run_all_tests()
    sys.exit(0 if success else 1)